import pptx
from pptx.util import Inches
from pptx.dml.color import RGBColor

# create presentation
prs = pptx.Presentation()

# set background color
background_color = RGBColor(255, 255, 255) # white color
prs.slide_width = Inches(13.3333)
prs.slide_height = Inches(7.5)
prs.slide_master.background.fill.solid()
prs.slide_master.background.fill.fore_color.rgb = background_color

# create title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = " UNIVERSITY"
subtitle = slide.placeholders[1]
subtitle.text = "Action Plan"

# create slides for university description
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "University Description"

bullet_slide_layout = prs.slide_layouts[1]
bullet_slide = prs.slides.add_slide(bullet_slide_layout)
shapes = bullet_slide.shapes

# add text to slide
shapes.title.text = "Worldwide VVIPs Required"
body_placeholder = shapes.placeholders[1]
tf = body_placeholder.text_frame

# add bullet points
tf.text = "Types of persons required:"
p = tf.add_paragraph()
p.text = "1. Accounting"
p.level = 1
p = tf.add_paragraph()
p.text = "2. Organic gardener"
p.level = 1
p = tf.add_paragraph()
p.text = "3. Clinical nutritionist"
p.level = 1
p = tf.add_paragraph()
p.text = "4. Catering and hotel management"
p.level = 1
p = tf.add_paragraph()
p.text = "5. Housekeeping"
p.level = 1
p = tf.add_paragraph()
p.text = "6. Geriatric qualified to take care of their elders"
p.level = 1
p = tf.add_paragraph()
p.text = "7. Well-behaving potential drivers and helpers"
p.level = 1
p = tf.add_paragraph()
p.text = "8. Loyal legal advisors"
p.level = 1

# create slides for programs offered
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Programs Offered"

bullet_slide_layout = prs.slide_layouts[1]
bullet_slide = prs.slides.add_slide(bullet_slide_layout)
shapes = bullet_slide.shapes

# add text to slide
shapes.title.text = "Programs Offered"
body_placeholder = shapes.placeholders[1]
tf = body_placeholder.text_frame

# add bullet points
tf.text = "1. Law with 1000 intake"
p = tf.add_paragraph()
p.text = "2. Mental health science courses clinical and non-clinical"
p.level = 1
p = tf.add_paragraph()
p.text = "3. Food science both clinical and non-clinical"
p.level = 1
p = tf.add_paragraph()
p.text = "4. MBA with all specializations"
p.level = 1
p = tf.add_paragraph()
p.text = "5. Nursing and geriatric nursing"
p.level = 1
p = tf.add_paragraph()
p.text = "6. Fine arts BFA and MFA"
p.level = 1
p = tf.add_paragraph()
p.text = "7. Intelligence courses to substitute raw and NIA"
p.level = 1

# add animations
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        for i, paragraph in enumerate(text_frame.paragraphs):
            for j, run in enumerate(paragraph.runs):
                if "bullet" in run.text:
                    text_range = run._r
                    bullet_size = 400000
                    x, y = text_range.left, text_range.top
                    prs.slides[0].shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        x-bullet_size,
                        y+bullet_size,
                        bullet_size/2,
                        bullet_size/2,
                    )
                    prs.slides[0].shapes.add_textbox(
                        x-bullet_size*3,
                        y+bullet_size,
                        bullet_size*2,
                        bullet_size/2,
                    ).text = "•"
                    paragraph.text = run.text.replace("bullet", "")

# save the presentation
prs.save("University Presentation.pptx")


